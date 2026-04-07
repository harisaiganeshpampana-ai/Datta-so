from flask import Flask, render_template, request, jsonify, send_file
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PyPDF2 import PdfReader
import io
import time
from docx2pdf import convert
import tempfile

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max

# Create uploads directory
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Popular student PDF resources
PDF_SOURCES = {
    "NCERT Books": "https://ncert.nic.in/textbook.php",
    "Previous Year Papers": "https://cbseacademic.nic.in/",
    "Sample Papers": "https://cbse.gov.in/",
    "Study Notes": "https://www.vedantu.com/"
}

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF"""
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except:
        text = "Unable to extract text from PDF"
    return text

def create_ppt_from_text(text, filename="presentation"):
    """Generate PPT from extracted text"""
    prs = Presentation()
    
    # Split text into chunks for slides
    paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
    slide_content = []
    
    # Group content for slides (max 8 lines per slide)
    current_slide = []
    for para in paragraphs:
        if len(current_slide) >= 8 or len(' '.join(current_slide)) > 800:
            if current_slide:
                slide_content.append('\n'.join(current_slide))
                current_slide = []
        current_slide.append(para)
    if current_slide:
        slide_content.append('\n'.join(current_slide))
    
    # Create slides
    for i, content in enumerate(slide_content, 1):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"Slide {i} - Study Notes"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = content
        
        # Format text
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(14)
    
    # Save PPT
    ppt_path = f"uploads/{filename}.pptx"
    prs.save(ppt_path)
    return ppt_path

@app.route('/')
def index():
    return render_template('index.html', sources=PDF_SOURCES)

@app.route('/upload', methods=['POST'])
def upload_pdf():
    if 'pdf' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['pdf']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and file.filename.lower().endswith('.pdf'):
        filename = os.path.splitext(file.filename)[0]
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)
        
        start_time = time.time()
        
        # Extract text
        text = extract_text_from_pdf(filepath)
        
        # Generate PPT
        ppt_path = create_ppt_from_text(text, filename)
        
        end_time = time.time()
        processing_time = round((end_time - start_time) * 1000, 2)  # ms
        
        return jsonify({
            'success': True,
            'ppt_url': f'/download/{os.path.basename(ppt_path)}',
            'processing_time': f'{processing_time}ms',
            'slides_count': len(text.split('\n\n'))
        })
    
    return jsonify({'error': 'Invalid file format'}), 400

@app.route('/download/<filename>')
def download_ppt(filename):
    return send_file(f'uploads/{filename}', as_attachment=True)

@app.route('/api/pdfs')
def get_pdf_sources():
    return jsonify(PDF_SOURCES)

if __name__ == '__main__':
    app.run(debug=True, port=5000)
